"""This script extract text from various format document to make llm fine tunning training data
"""
import argparse
import json
import os
import pprint
import pathlib
import re

import pandas as pd
import pdfplumber
from langchain.document_loaders import TextLoader
from pptx import Presentation
from docx import Document
from docx.document import Document as _Document
from docx.oxml.text.paragraph import CT_P
from docx.oxml.table import CT_Tbl
from docx.table import _Cell, Table, _Row
from docx.text.paragraph import Paragraph
from HwpParser import HWPExtractor

class TextExtract:
    """
    Document extration main class
    """
    def __init__(self, bucket_name=None):
        self.del_table = True
        self.cvt_image = False

    def is_inside_any_table(self, word_bbox, tables):
        """단어가 테이블 내부에 있는지 확인"""
        word_x0, word_y0, word_x1, word_y1 = word_bbox
        word_center_y = (word_y0 + word_y1) / 2
        word_center_x = (word_x0 + word_x1) / 2

        for table in tables:
            table_x0, table_y0, table_x1, table_y1 = table.bbox
            if (table_x0 <= word_center_x <= table_x1 and
                    table_y0 <= word_center_y <= table_y1):
                return True
        return False

    def get_contexttable_pdffile_by_plumber(self, source_file_name: str) -> list:
        result = {
            "doc_meta": [],
        }
        try:
            with pdfplumber.open(source_file_name) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    # 페이지의 모든 요소를 담을 리스트
                    page_elements = []

                    # 테이블 찾기
                    tables = page.find_tables()

                    # 테이블 처리
                    for table_num, table in enumerate(tables, 1):
                        df = pd.DataFrame(table.extract())
                        if len(df) > 0:  # 빈 테이블 제외
                            df.columns = df.iloc[0]
                            df = df.iloc[1:]
                            markdown_table = df.to_markdown(index=False)

                            # 테이블의 중심 y 좌표 계산
                            y_center = (table.bbox[1] + table.bbox[3]) / 2

                            page_elements.append({
                                'type': 'table',
                                'context': markdown_table,
                                'y_center': y_center,
                                'page': page_num,
                                'table_number': table_num,
                                'bbox': table.bbox
                            })

                    # 텍스트 추출 (테이블 영역 제외)
                    words = page.extract_words(keep_blank_chars=True, x_tolerance=3, y_tolerance=3)

                    # 테이블 영역 외의 단어들만 모으기
                    current_line = []
                    current_line_number = 1
                    current_y = None

                    for word in words:
                        word_bbox = (word['x0'], word['top'], word['x1'], word['bottom'])
                        if self.is_inside_any_table(word_bbox, tables):
                            continue
                        word_y_center = (word['top'] + word['bottom']) / 2
                        if current_y is None:
                            current_y = word_y_center

                        if abs(word_y_center - current_y) < 5:
                            current_line.append(word['text'])
                        else:
                            # 새로운 라인 시작
                            if current_line:
                                line_text = " ".join(current_line)
                                if line_text.strip():  # 빈 라인 제외
                                    page_elements.append({
                                        'type': 'text',
                                        'context': f"{line_text}",
                                        'y_center': current_y,
                                        'page': page_num,
                                        'line': current_line_number
                                    })
                                    current_line_number += 1
                            current_line = [word['text']]
                            current_y = word_y_center

                    # 마지막 라인 처리
                    if current_line:
                        line_text = " ".join(current_line)
                        if line_text.strip():
                            page_elements.append({
                                'type': 'text',
                                'context': f"{line_text}",
                                'y_center': current_y,
                                'page': page_num,
                                'line': current_line_number
                            })

                    # y 위치에 따라 요소들을 정렬
                    page_elements.sort(key=lambda x: x['y_center'])

                    # 텍스트 요소들을 청크로 분할하며 결과에 추가
                    current_text = ""
                    current_text_elements = []

                    for element in page_elements:
                        if element['type'] == 'text':
                            current_text += element['context'] + "\n"
                            current_text_elements.append({
                                'page': element['page'],
                                'line': element['line']
                            })
                            start_ln = current_text_elements[0]['line']
                            fist_end_page_line_num = {'start_line': start_ln}
                            result["doc_meta"].append({
                                "type": "text",
                                "page": page_num,
                                "context": current_text,
                                "line_pos": fist_end_page_line_num
                            })
                            current_text = ""
                            current_text_elements = []
                        else:  # 테이블인 경우
                            # 남은 텍스트가 있으면 먼저 처리
                            if current_text:
                                result["doc_meta"].append({
                                    "type": "text",
                                    "page": page_num,
                                    "context": current_text,
                                    "line_pos": fist_end_page_line_num
                                })
                                current_text = ""
                                current_text_elements = []

                            # 테이블 추가
                            result["doc_meta"].append({
                                "type": "table",
                                "page": element['page'],
                                "table_number": element['table_number'],
                                "context": element['context']
                            })

                    # 페이지의 마지막 텍스트 처리
                    if current_text:
                        start_ln = current_text_elements[0]['line']
                        fist_end_page_line_num = {'start_line': start_ln}
                        result["context"].append({
                            "type": "text",
                            "page": page_num,
                            "context": current_text,
                            "line_pos": fist_end_page_line_num
                        })
                return None, result
        except Exception as e:
            print(str(e))
            print("No next doc")
            return None, None


    def iter_doc_blocks(self, parent):
        """
        yield document element type
        :param parent: word paragraph object
        :return: word element value
        """
        if isinstance(parent, _Document):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        elif isinstance(parent, _Row):
            parent_elm = parent._tr
        else:
            raise ValueError("something's not right")
        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    def get_msword_content(self, file_path):
        """
            MS Word 문서에서 텍스트를 추출하고 페이지 번호와 줄 번호를 포함하여 반환합니다.

            Args:
                file_path (str): Word 문서의 경로

            Returns:
                list: 각 줄의 정보를 담은 딕셔너리 리스트
                    - page_number: 페이지 번호
                    - line_number: 줄 번호
                    - content: 텍스트 내용
            """
        try:
            doc = Document(file_path)
        except Exception as e:
            raise Exception(f"문서를 열 수 없습니다: {str(e)}")
        current_page = 1
        line_number = 1
        doc_meta = []
        whole_page_extractinfo = ''
        for paragraph in doc.paragraphs:
            # 페이지 나누기 확인
            if paragraph._element.xpath('.//w:lastRenderedPageBreak'):
                current_page += 1
            text_list = paragraph.text.split("\n")
            if len(text_list) > 0:  # 빈 줄 제외
                for text in text_list:
                    whole_page_extractinfo += text + '\n'
                    doc_meta.append({"type": "text",
                                     "page": current_page,
                                     "context": text,
                                     "line_pos": {"start_line": line_number}})
                    line_number += 1
        return whole_page_extractinfo, doc_meta


    def get_ppt_context(self, ppt_prs):
        """
         extraction pptx file contents from ppt object
        :param ppt_prs: ppt object
        :return:  extracted ppt contents
        """
        pptx_contents = ''
        for _, slide in enumerate(ppt_prs.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        pptx_contents += run.text + '\r\n'
        return pptx_contents



    def extract_file_content(self, file_path):
        """
        extraction contents from various file format
        :param file_name: file name(object name) in object storage
        :return: extract text
        """
        formed_clear_contents = ''
        f_extension = pathlib.Path(file_path).suffix
        f_extension = f_extension.lower()

        if f_extension.endswith('.pdf'):
            _, infor_meta = self.get_contexttable_pdffile_by_plumber(file_path)
            if infor_meta is not None:
                self.preprocess_meta = {'origin_path':file_path,
                                            'doc_meta':infor_meta['doc_meta']}
            else:
                self.preprocess_meta = None

        elif f_extension.endswith('.hwp'):
            hwp_obj = HWPExtractor(file_path)
            hwp_text = hwp_obj.get_text()
            formed_clear_contents = hwp_text

        elif f_extension.endswith('.docx') or f_extension.endswith('.doc'):
            formed_clear_contents, doc_meta = self.get_msword_content(file_path)
            self.preprocess_meta.append({'origin_path': file_path,
                                        'doc_meta': doc_meta})

        elif f_extension.endswith('.txt'):
            loader = TextLoader(file_path)
            docs = loader.load()
            for page in docs:
                formed_clear_contents += page.page_content

        elif f_extension.endswith('.xlsx') or f_extension.endswith('.xls'):
            df = pd.read_excel(file_path)
            df_markdown = df.to_markdown()
            formed_clear_contents = df_markdown

        elif f_extension.endswith('.csv'):
            df = pd.read_csv(file_path)
            df_markdown = df.to_markdown()
            formed_clear_contents = df_markdown

        elif f_extension.endswith('.pptx') or f_extension.endswith('.ppt'):
            try:
                prs = Presentation(file_path)
                formed_clear_contents = self.get_ppt_context(prs)
                print(formed_clear_contents)
            except Exception as e:
                print(e)
        else:
            print("Error: invalid file type")
            print(file_path)
        return 1, formed_clear_contents, "ok"

    def is_pdf(self, file_path):
        # PDF 파일은 '%PDF-' 로 시작합니다
        try:
            with open(file_path, 'rb') as file:
                header = file.read(5)
                return header.startswith(b'%PDF-')
        except Exception as e:
            print(f"Error: {e}")
            return False

    def extract_from_src_doc(self, path: str):
        """
        extract contents from all object in target bucket
        :return:
        """
        # types = ('.pdf', '.docx')
        types = ('.pdf')
        src_meta_path = '../meta_dumps'
        if os.path.exists(src_meta_path) == False:
            os.makedirs(src_meta_path, exist_ok=True)

        origin_file_list = []
        for metafile in pathlib.Path(src_meta_path).rglob("*.json"):
            with open(metafile, "r", encoding="utf-8") as file:
                print(metafile)
                file_info = json.load(file)
                orgin_file = file_info['origin_path']
                origin_file_list.append(orgin_file)

        for idx, filepath in enumerate(pathlib.Path(path).rglob('*.*')):
            if filepath.is_file():
                sfile_path = str(filepath)
                if sfile_path.endswith(types):
                    if sfile_path in origin_file_list: continue
                    print(f"In process {sfile_path}")
                    if self.is_pdf(sfile_path) == False: continue

                    self.extract_file_content(sfile_path)
                    with open(f'../meta_dumps/{idx}_metadump.json', "w", encoding='utf-8') as fp:
                        if self.preprocess_meta is not None:
                            fp.write(json.dumps(self.preprocess_meta, ensure_ascii=False))
                        else:
                            fp.write(json.dumps({"origin_path": sfile_path,
                                                      "status":"Error"},
                                                      ensure_ascii=False))

    def dump_data(self, meta_path):
        for metafile in pathlib.Path(meta_path).rglob("*.json"):
            with open(metafile, "r", encoding="utf-8") as file:
                file_info = json.load(file)
                pprint.pprint(file_info)

te = TextExtract()
te.del_table = True

parser = argparse.ArgumentParser()
parser.add_argument('-input', help='추출할 파일이 있는 경로를 넣어주세요')
args = parser.parse_args()
if __name__ == '__main__' :
    print(f'추출 대상 경로: ', args.input)
    te.extract_from_src_doc(path=args.input)
    te.dump_data("../meta_dumps")

